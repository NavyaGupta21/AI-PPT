import streamlit as st
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
from langchain_core.messages import HumanMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import PromptTemplate
import requests
from googlesearch import search
from urllib.parse import urljoin
from io import BytesIO

llm = ChatGoogleGenerativeAI(
    model="gemini-2.5-flash",
    temperature=0.7,
    google_api_key=st.secrets["GOOGLE_API_KEY"],
    max_tokens=4000,
)

st.title("PPT Generator")

url_mode = st.radio("URL Input Mode", ("Manual", "Google Search"))

def scrap(urls, no_of_slides):
    try:
        combined_text = ""
        images_urls = []
        st.info("Scraping content...")

        for url in urls:
            try:
                response = requests.get(url, timeout=10)
                response.raise_for_status()
                soup = BeautifulSoup(response.text, "html.parser")

                text_content = []
                for tag in ["h1", "h2", "h3", "p", "li", "article", "section", "div"]:
                    elements = soup.find_all(tag)
                    text_content.extend([el.get_text().strip() for el in elements if el.get_text().strip()])

                if not text_content:
                    for script in soup(["script", "style"]):
                        script.extract()
                    text_content = [line.strip() for line in soup.get_text().splitlines() if line.strip()]

                chars_per_slide = 300
                combined_text += " ".join(text_content)[:no_of_slides * chars_per_slide] + "\n\n"

                valid_image_extensions = (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg")
                current_page_images = [
                    urljoin(url, img['src']) for img in soup.find_all("img", src=True)
                    if img['src'].strip() and
                       (img['src'].startswith("http") or img['src'].startswith("/")) and 
                       any(ext in img['src'].lower() for ext in valid_image_extensions) and
                       "Special:" not in img['src'] and 
                       "type=1x1" not in img['src'] 
                ]
                images_urls.extend(current_page_images)
                
            except requests.exceptions.RequestException as req_err:
                st.error(f"Error fetching {url}: {req_err}")
                continue
            except Exception as parse_err:
                st.error(f"Error parsing content from {url}: {parse_err}")
                continue

        if not combined_text.strip():
            st.warning("No content available to generate slides")

        st.info("Content scraped")

        prompt = PromptTemplate(
            input_variables=["content", "images", "slides_count"],
            template="""
            Based on the following scraped content {content} and images {images}, generate {slides_count} PowerPoint slides.

            Each slide should include:
            - A clear and descriptive Title.
            - Bullet points summarizing key details from the content.
            - Include only those images from the provided list of images ({images}) that directly align with and enhance the understanding of the content ({content}). Avoid adding irrelevant or decorative images.
            - Ensure the content is concise, directly related to the provided text, and logically structured.
            - Format each slide clearly, starting with "Slide [Number]: [Title]" followed by bullet points and any related images(each slide should have only one image).

            # Example:
            Slide 1: Introduction to Topic
            - Key point 1
            - Key point 2
            http://example.com/image1.jpg

            Slide 2: Key Concept A
            - Detail 1
            - Detail 2
            http://example.com/image2.jpg"

            Follow the given format only

            Content to use for slide generation:
            {content}

            Images to use for slide generation:
            {images}
            """
        )
        query = prompt.format(content=combined_text, images=images_urls, slides_count=no_of_slides)
        llm_response = llm.invoke([HumanMessage(content=query)])

        response_content = (
            llm_response.content if hasattr(llm_response, "content") else str(llm_response)
        )

        slide_sections = response_content.split("Slide ")
        refined_content = []

        for section in slide_sections:
            section = section.strip()
            if not section:
                continue

            first_idx = section.find('\n')
            if first_idx == -1:
                title = section.split(':', 1)[-1].strip()
                points_text = ""
                images_url = ""
            else:
                title_line = section[:first_idx]
                title = title_line.split(':', 1)[-1].strip()
                remaining_content = section[first_idx+1:].strip()
                
                second_idx = remaining_content.find("http")
                if second_idx != -1:
                    points_text = remaining_content[:second_idx].strip()
                    images_url = remaining_content[second_idx:].strip()
                else:
                    points_text = remaining_content.strip()
                    images_url = ""

            if points_text or images_url:
                refined_content.append({"title": title, "points": points_text, "images": images_url})

        """ if refined_content:
            for i, slide in enumerate(refined_content):
                st.subheader(f"Slide {i + 1}: {slide['title']}")
                st.markdown(
                    "\n".join([f"- {line.strip()}" if not line.strip().startswith(("-", "*")) else line.strip()
                               for line in slide["points"].split('\n') if line.strip()])
                )
                if slide["images"] and slide["images"] != '.':
                    st.image(slide["images"], caption=f"Image for Slide {i+1}")
                else:
                    st.warning(f"No images available for Slide {i+1}.") """

        template_path = "Bracket design.pptx"
        ppt = Presentation(template_path)
        for i, slide_data_dict in enumerate(refined_content):
            try:
                slide_layout = ppt.slide_layouts[17]
                slide = ppt.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data_dict['title']
                else:
                    st.warning(f"Slide {i+1}: No title placeholder found in selected layout.")
        
               
                content_placeholder = None
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == 14:
                        content_placeholder = placeholder
                        break
        
                if content_placeholder and content_placeholder.has_text_frame:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    for point in slide_data_dict["points"].split('\n'):
                        if point.strip():
                            p = text_frame.add_paragraph()
                            p.text = point
                else:
                    st.warning(f"Slide {i+1}: No suitable content placeholder for bullet points in selected layout.") 
        
                if slide_data_dict["images"] and slide_data_dict["images"] != '.':
                    try:
                        img_response = requests.get(slide_data_dict["images"])
                        img_response.raise_for_status()
                        image_stream = BytesIO(img_response.content)
                        
                        picture_placeholder = None
                        for placeholder in slide.placeholders:
                            if placeholder.placeholder_format.idx == 13:
                                picture_placeholder = placeholder
                                break
                                
                        if picture_placeholder:
                            slide.shapes.add_picture(image_stream, 
                                             picture_placeholder.left, 
                                             picture_placeholder.top, 
                                             width=picture_placeholder.width, 
                                             height=picture_placeholder.height)
                        else:
                            st.warning(f"Slide {i+1}: No Picture Placeholder found in selected layout.")
                    
                    except requests.exceptions.RequestException as img_req_err:
                         st.warning(f"Could not download image {slide_data_dict['images']} for Slide {i+1}: {img_req_err}")
                    except Exception as img_add_err:
                         st.warning(f"Error adding image to Slide {i+1}: {img_add_err}")
                else:
                    st.warning(f"No images provided for Slide {i+1}.")
            except Exception as slide_creation_err:
                st.error(f"Error creating Slide {i+1}: {slide_creation_err}")
                continue
        ppt_io = BytesIO()
        ppt.save(ppt_io)
        ppt_io.seek(0)

        file_path = "PPT.pptx"
        with open(file_path, "wb") as f:
            f.write(ppt_io.getbuffer())
            st.info(f"PowerPoint file saved at {file_path}")
  
    except Exception as e:
        st.error(f"An unexpected error occurred during processing: {e}")
    
urls = []
if url_mode == "Manual":
    no_of_url = st.number_input("No of URLs", min_value=1, max_value=5, value=1)
    for i in range(no_of_url):
        url = st.text_input(f"Enter URL {i + 1}", key=f"url_{i}")
        urls.append(url)
    no_of_slides = st.number_input("Slides to Generate", min_value=1, max_value=20, value=5)
    if st.button("Generate"):
        scrap(urls,no_of_slides)
    
else:
    search_query = st.text_input("Search Query")
    required_valid_urls = st.number_input("No. of URLs", min_value=1, max_value=10, value=3)
    no_of_slides = st.number_input("Slides to Generate", min_value=1, max_value=20, value=5)
    if st.button("Generate"):
        urls = []
        valid_urls = []
        st.info("Fetching URLs...")
        search_results = search(search_query, num_results=20)
        for url in search_results:
            if len(valid_urls) >= required_valid_urls:
                break
            else:
                try:
                    response = requests.head(url, timeout=5)
                    if response.status_code == 200:
                        valid_urls.append(url)
                except Exception:
                    continue
                    
        if len(valid_urls) == required_valid_urls:
            urls = valid_urls
            st.success("Successfully fetched URLs:")
            for i, url in enumerate(urls):
                st.write(f"{i + 1}. {url}")
            scrap(urls,no_of_slides)
        else:
            st.warning(f"Only {len(valid_urls)} valid URLs found. Please refine your search.")   